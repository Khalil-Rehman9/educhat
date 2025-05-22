"""
PPTX document processor using python-pptx for text extraction.
"""

import os
import sys
from datetime import datetime
from typing import Dict, Any, List, Optional
import pptx

# Add parent directory to path for imports
sys.path.append(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

from app.backend.document_processors.base_processor import BaseDocumentProcessor

class PPTXProcessor(BaseDocumentProcessor):
    def validate_doc_format(self):
        if not self.file_path.lower().endswith('.pptx'):
            raise ValueError("Invalid file extension for PPTXProcessor.")
    """
    Processor for Microsoft PowerPoint (PPTX) presentations using python-pptx.
    """
    
    def extract_text(self) -> str:
        self.validate_doc_format()
        """
        Extract text content from a PPTX presentation.
        
        Returns:
            The extracted text content as a string
        """
        try:
            text_content = ""
            
            # Open the PPTX file
            presentation = pptx.Presentation(self.file_path)
            
            # Extract text from each slide
            for i, slide in enumerate(presentation.slides):
                slide_num = i + 1
                text_content += f"\n--- Slide {slide_num} ---\n"
                
                # Get slide title if it exists
                if slide.shapes.title:
                    text_content += f"Title: {slide.shapes.title.text}\n\n"
                
                # Get text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Skip the title that we've already added
                        if shape == slide.shapes.title:
                            continue
                        text_content += shape.text + "\n"
                
                # Add a separator between slides
                text_content += "\n"
            
            return text_content.strip()
        
        except Exception as e:
            error_msg = f"Error extracting text from PPTX: {str(e)}"
            print(error_msg)
            return f"ERROR: {error_msg}"
    
    def extract_metadata(self) -> Dict[str, Any]:
        """
        Extract metadata from a PPTX presentation.
        
        Returns:
            Dictionary of metadata
        """
        metadata = {
            "file_path": self.file_path,
            "document_id": self.document_id,
            "processor": "PPTXProcessor",
            "extraction_date": datetime.now().isoformat(),
            "slide_count": 0,
            "title": None,
            "author": None,
            "subject": None,
            "keywords": None,
            "comments": None,
            "last_modified_by": None,
            "revision": None,
            "created": None,
            "modified": None
        }
        
        try:
            # Open the PPTX file
            presentation = pptx.Presentation(self.file_path)
            
            # Get slide count
            metadata["slide_count"] = len(presentation.slides)
            
            # Extract document properties
            core_props = presentation.core_properties
            
            if core_props:
                if hasattr(core_props, 'title') and core_props.title:
                    metadata["title"] = core_props.title
                if hasattr(core_props, 'author') and core_props.author:
                    metadata["author"] = core_props.author
                if hasattr(core_props, 'subject') and core_props.subject:
                    metadata["subject"] = core_props.subject
                if hasattr(core_props, 'keywords') and core_props.keywords:
                    metadata["keywords"] = core_props.keywords
                if hasattr(core_props, 'comments') and core_props.comments:
                    metadata["comments"] = core_props.comments
                if hasattr(core_props, 'last_modified_by') and core_props.last_modified_by:
                    metadata["last_modified_by"] = core_props.last_modified_by
                if hasattr(core_props, 'revision') and core_props.revision:
                    metadata["revision"] = core_props.revision
                if hasattr(core_props, 'created') and core_props.created:
                    metadata["created"] = core_props.created.isoformat() if hasattr(core_props.created, 'isoformat') else str(core_props.created)
                if hasattr(core_props, 'modified') and core_props.modified:
                    metadata["modified"] = core_props.modified.isoformat() if hasattr(core_props.modified, 'isoformat') else str(core_props.modified)
        
        except Exception as e:
            metadata["error"] = str(e)
        
        return metadata
    
    def extract_slides(self) -> List[Dict[str, Any]]:
        """
        Extract detailed slide information.
        
        Returns:
            List of dictionaries containing slide data
        """
        slides_data = []
        
        try:
            # Open the PPTX file
            presentation = pptx.Presentation(self.file_path)
            
            # Process each slide
            for i, slide in enumerate(presentation.slides):
                slide_num = i + 1
                
                # Get slide title
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text
                
                # Get all text from the slide
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Skip the title that we've already captured
                        if shape == slide.shapes.title:
                            continue
                        texts.append(shape.text)
                
                # Create slide data entry
                slide_data = {
                    "slide_number": slide_num,
                    "title": title,
                    "content": "\n".join(texts),
                    "has_notes": hasattr(slide, "notes_slide") and slide.notes_slide and slide.notes_slide.notes_text_frame.text.strip() != "",
                    "notes": slide.notes_slide.notes_text_frame.text if hasattr(slide, "notes_slide") and slide.notes_slide else ""
                }
                
                slides_data.append(slide_data)
        
        except Exception as e:
            print(f"Error extracting slides: {str(e)}")
        
        return slides_data 